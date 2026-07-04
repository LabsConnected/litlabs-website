"""
Small generator to create DOCX, XLSX and PPTX from project's markdown sources.
Run: python3 scripts/generate_office_files.py
"""
from pathlib import Path
try:
    from docx import Document
    from openpyxl import Workbook
    from pptx import Presentation
except Exception as e:
    print('Missing dependencies:', e)
    raise

# Paths
repo_root = Path('.')
md_report = repo_root / 'tasks' / 'completed' / 'studio-health-report-2026-06-08.md'
md_prd = repo_root / 'prds' / 'ai-studio-optimization.md'
outputs = repo_root / 'outputs'
outputs.mkdir(exist_ok=True)

# Create DOCX from report
if md_report.exists():
    text = md_report.read_text(encoding='utf-8')
    doc = Document()
    doc.add_heading('Studio Health Report — 2026-06-08', level=1)
    doc.add_paragraph('TL;DR: Converted from markdown source.')
    for line in text.splitlines():
        if line.startswith('# '):
            continue
        if line.startswith('## '):
            doc.add_heading(line[3:].strip(), level=2)
        elif line.startswith('---'):
            doc.add_page_break()
        else:
            doc.add_paragraph(line)
    doc.save(outputs / 'studio-health-report-2026-06-08.docx')

# Create XLSX with action items
wb = Workbook()
ws = wb.active
ws.title = 'Action Items'
rows = [
    ('#','Action','Effort','Status','Owner'),
    (1,'Add global Agent Chat chip to Navbar','5 min','Optional','Frontend'),
    (2,'Wire /api/wallet to studio top bar coin badge','10 min','Optional','Frontend'),
    (3,'Make ModelBadge reactive to provider','30 min','Optional','Frontend'),
    (4,'Patch brain.sh + monitor.sh to use cloudflared-main','2 min','Recommended','Devops'),
    (5,'Add systemd timer for agent-monitor.service','5 min','Optional','Devops'),
    (6,'Verify OpenRouter key + run /api/llm/health','1 min','Recommended','Backend'),
]
for r in rows:
    ws.append(r)
# Add verification sheet
ws2 = wb.create_sheet('Verification Commands')
ws2.append(('command',))
ws2.append(("wsl -e bash -c \"systemctl status cloudflared-main --no-pager\"",))
ws2.append(("curl -s http://localhost:3000/api/llm/health | python -m json.tool",))
wb.save(outputs / 'studio-health-actions-2026-06-08.xlsx')

# Create PPTX from PRD headings
prs = Presentation()
if md_prd.exists():
    text = md_prd.read_text(encoding='utf-8')
    # Title slide
    s = prs.slides.add_slide(prs.slide_layouts[0])
    title = s.shapes.title
    subtitle = s.placeholders[1]
    title.text = 'AI Studio & Performance Optimization'
    subtitle.text = 'TL;DR — consolidated PRD'
    # Simple parsing: split by H2 headings
    sections = []
    cur = None
    for line in text.splitlines():
        if line.startswith('## '):
            if cur:
                sections.append(cur)
            cur = {'title': line[3:].strip(), 'body': []}
        elif cur is not None:
            cur['body'].append(line)
    if cur:
        sections.append(cur)
    # Add up to 4 more slides
    for sec in sections[:4]:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = sec['title']
        tf = slide.shapes.placeholders[1].text_frame
        for p in sec['body'][:6]:
            if p.strip():
                p_obj = tf.add_paragraph()
                p_obj.text = p.strip()
prs.save(outputs / 'ai-studio-prd-2026-07-02.pptx')

print('Generated outputs in outputs/')
