from openpyxl import load_workbook
from pptx import Presentation
from docx import Document
from pathlib import Path

out = Path('outputs')

# Verify XLSX
wb = load_workbook(out / 'studio-health-actions-2026-06-08.xlsx')
print('XLSX sheets:', wb.sheetnames)
for ws in wb.worksheets:
    print(f'  {ws.title}: {ws.max_row} rows x {ws.max_column} cols')
    for row in ws.iter_rows(values_only=True):
        print('   ', row)

# Verify PPTX
prs = Presentation(out / 'ai-studio-prd-2026-07-02.pptx')
print('\nPPTX slides:', len(prs.slides))
for i, slide in enumerate(prs.slides, 1):
    title = slide.shapes.title.text if slide.shapes.title else 'NO TITLE'
    print(f'  Slide {i}: {title}')

# Verify DOCX
doc = Document(out / 'studio-health-report-2026-06-08.docx')
print('\nDOCX paragraphs:', len(doc.paragraphs))
print('  First heading:', doc.paragraphs[1].text)
